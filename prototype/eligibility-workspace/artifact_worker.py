"""Render approved Studio projects into real editable artifacts."""

from __future__ import annotations

import base64
import hashlib
import html
import json
import pathlib
import re
import shutil
import subprocess
import uuid

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


BLUE = RGBColor(22, 50, 79)
ORANGE = RGBColor(229, 104, 42)
SLATE = RGBColor(69, 82, 96)
LIGHT = RGBColor(240, 244, 247)
PPT_BLUE = PptRGBColor(22, 50, 79)
PPT_ORANGE = PptRGBColor(229, 104, 42)
PPT_WHITE = PptRGBColor(255, 255, 255)
PPT_SLATE = PptRGBColor(69, 82, 96)


def _safe_slug(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return slug[:80] or "artifact"


def _plain_text(value: str) -> str:
    """Remove common Markdown syntax while preserving the authored wording."""
    text = str(value or "")
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"(?m)^\s{0,3}#{1,6}\s+", "", text)
    text = re.sub(r"(?m)^\s*>\s?", "", text)
    text = re.sub(r"(?<=\w)\*\*(?=\w)", " ", text)
    text = re.sub(r"(?<=\w)__(?=\w)", " ", text)
    text = text.replace("**", "").replace("__", "")
    text = re.sub(r"(?<!\w)[*_](?=\S)|(?<=\S)[*_](?!\w)", "", text)
    return text.strip()


def _checksum(path: pathlib.Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return f"sha256:{digest.hexdigest()}"


def _file_record(path: pathlib.Path, render_id: str, fmt: str) -> dict:
    return {
        "file_id": f"file:{uuid.uuid4()}",
        "filename": path.name,
        "format": fmt,
        "size": path.stat().st_size,
        "checksum": _checksum(path),
        "download_url": f"/artifacts/{render_id}/{path.name}",
    }


def _decode_images(project: dict, output_dir: pathlib.Path) -> list[dict]:
    saved: list[dict] = []
    for index, slot in enumerate(project.get("image_slots", []), start=1):
        asset = slot.get("asset") or {}
        data_url = asset.get("data_url")
        if not data_url:
            continue
        match = re.match(r"^data:image/(png|jpeg|webp);base64,(.+)$", data_url, re.DOTALL)
        if not match:
            raise ValueError("Uploaded images must be PNG, JPEG, or WebP data URLs.")
        extension = "jpg" if match.group(1) == "jpeg" else match.group(1)
        path = output_dir / f"image-{index}.{extension}"
        path.write_bytes(base64.b64decode(match.group(2), validate=True))
        saved.append({**slot, "path": path})
    return saved


def _set_repeat_table_header(row) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), "true")
    tr_pr.append(tbl_header)


def _shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    tc_pr.append(shading)


def _style_doc(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    normal = document.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.font.color.rgb = SLATE
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.25
    for name, size, color, before, after in (
        ("Title", 26, BLUE, 0, 14),
        ("Heading 1", 18, BLUE, 14, 6),
        ("Heading 2", 14, ORANGE, 10, 4),
        ("Heading 3", 11, BLUE, 8, 3),
    ):
        style = document.styles[name]
        style.font.name = "Calibri"
        style.font.size = Pt(size)
        style.font.bold = name != "Normal"
        style.font.color.rgb = color
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)
        style.paragraph_format.keep_with_next = True


def _add_doc_metadata(document: Document, project: dict, brief: dict) -> None:
    table = document.add_table(rows=2, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    widths = [Inches(2.1), Inches(2.1), Inches(2.1)]
    values = [
        ("Audience", project.get("audience", "")),
        ("Brief version", str(project.get("brief_version", ""))),
        ("Template", project.get("template_id", "")),
        ("Objective", project.get("objective", "")),
        ("Sources", str(len(project.get("source_ids", [])))),
        ("Status", "Approved content input"),
    ]
    for index, cell in enumerate(table._cells):
        cell.width = widths[index % 3]
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        label, value = values[index]
        paragraph = cell.paragraphs[0]
        label_run = paragraph.add_run(f"{label}\n")
        label_run.bold = True
        label_run.font.color.rgb = BLUE
        paragraph.add_run(value)
        _shade_cell(cell, "F0F4F7" if index < 3 else "FFFFFF")


def _render_docx(project: dict, brief: dict, output_dir: pathlib.Path, images: list[dict]) -> pathlib.Path:
    document = Document()
    _style_doc(document)
    document.core_properties.title = project["title"]
    document.core_properties.subject = project["objective"]
    document.add_heading(project["title"], 0)
    _add_doc_metadata(document, project, brief)
    if project.get("summary"):
        document.add_heading("At a glance", level=1)
        document.add_paragraph(project["summary"])
    document.add_heading("Key guidance", level=1)
    points = brief.get("points") or []
    for index, point in enumerate(points, start=1):
        use = str(point.get("intended_use", "key_fact")).replace("_", " ").title()
        document.add_heading(f"{index}. {use}", level=2)
        document.add_paragraph(_plain_text(point.get("statement", "")))
        citations = point.get("citations") or []
        if citations:
            citation = document.add_paragraph()
            citation.style = document.styles["Caption"]
            run = citation.add_run("Source: " + "; ".join(item.get("label", "") for item in citations))
            run.italic = True
            run.font.color.rgb = SLATE
    if images:
        document.add_heading("Supporting visuals", level=1)
        for item in images:
            document.add_picture(str(item["path"]), width=Inches(6.25))
            caption = document.add_paragraph(item.get("caption") or item.get("alt_text") or "")
            caption.style = document.styles["Caption"]
            caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_heading("Source record", level=1)
    source_table = document.add_table(rows=1, cols=2)
    source_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    source_table.style = "Table Grid"
    source_table.rows[0].cells[0].text = "Source ID"
    source_table.rows[0].cells[1].text = "Role"
    _set_repeat_table_header(source_table.rows[0])
    for cell in source_table.rows[0].cells:
        _shade_cell(cell, "16324F")
        for run in cell.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.bold = True
    for source_id in project.get("source_ids", []):
        cells = source_table.add_row().cells
        cells[0].text = source_id
        cells[1].text = "Approved brief lineage"
    path = output_dir / f"{_safe_slug(project['title'])}.docx"
    document.save(path)
    return path


def _convert_pdf(docx_path: pathlib.Path, output_dir: pathlib.Path) -> pathlib.Path | None:
    office = shutil.which("soffice") or shutil.which("libreoffice")
    if not office:
        return None
    result = subprocess.run(
        [office, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(docx_path)],
        check=False,
        capture_output=True,
        text=True,
        timeout=120,
    )
    pdf_path = output_dir / f"{docx_path.stem}.pdf"
    if result.returncode != 0 or not pdf_path.exists():
        raise ValueError(f"PDF conversion failed: {(result.stderr or result.stdout).strip()[:300]}")
    return pdf_path


def _add_slide_title(slide, title: str, kicker: str = "") -> None:
    if kicker:
        box = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.35), PptInches(11.6), PptInches(0.35))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = kicker.upper()
        run.font.size = PptPt(10)
        run.font.bold = True
        run.font.color.rgb = PPT_ORANGE
    title_box = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.75), PptInches(11.7), PptInches(0.8))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = PptPt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = PPT_BLUE


def _add_footer(slide, project: dict, number: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.7), PptInches(7.02), PptInches(11.9), PptInches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = PPT_ORANGE
    line.line.fill.background()
    box = slide.shapes.add_textbox(PptInches(0.72), PptInches(7.08), PptInches(11.8), PptInches(0.25))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = f"{project.get('title', '')}   •   Brief v{project.get('brief_version', '')}   •   {number}"
    paragraph.font.size = PptPt(8)
    paragraph.font.color.rgb = PPT_SLATE


def _render_pptx(project: dict, brief: dict, output_dir: pathlib.Path, images: list[dict]) -> pathlib.Path:
    deck = Presentation()
    deck.slide_width = PptInches(13.333333)
    deck.slide_height = PptInches(7.5)
    blank = deck.slide_layouts[6]
    slide = deck.slides.add_slide(blank)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PPT_BLUE
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.22), PptInches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = PPT_ORANGE; accent.line.fill.background()
    title_box = slide.shapes.add_textbox(PptInches(0.85), PptInches(1.55), PptInches(10.9), PptInches(1.6))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = project["title"]
    paragraph.font.size = PptPt(36); paragraph.font.bold = True; paragraph.font.color.rgb = PPT_WHITE
    sub = slide.shapes.add_textbox(PptInches(0.88), PptInches(3.45), PptInches(9.8), PptInches(1.1))
    sub_p = sub.text_frame.paragraphs[0]
    sub_p.text = project["objective"]
    sub_p.font.size = PptPt(20); sub_p.font.color.rgb = PPT_WHITE
    audience = slide.shapes.add_textbox(PptInches(0.88), PptInches(6.35), PptInches(8.5), PptInches(0.4))
    audience_p = audience.text_frame.paragraphs[0]
    audience_p.text = project["audience"]
    audience_p.font.size = PptPt(11); audience_p.font.color.rgb = PPT_WHITE
    points = brief.get("points") or []
    for index, point in enumerate(points, start=1):
        slide = deck.slides.add_slide(blank)
        _add_slide_title(slide, point.get("intended_use", "Key point").replace("_", " ").title(), f"Approved point {index}")
        body = slide.shapes.add_textbox(PptInches(0.9), PptInches(1.85), PptInches(8.25), PptInches(3.9))
        body_tf = body.text_frame
        body_tf.word_wrap = True
        body_p = body_tf.paragraphs[0]
        body_p.text = _plain_text(point.get("statement", ""))
        body_p.font.size = PptPt(23); body_p.font.color.rgb = PPT_BLUE
        citations = point.get("citations") or []
        if citations:
            citation = slide.shapes.add_textbox(PptInches(0.9), PptInches(5.95), PptInches(8.4), PptInches(0.62))
            citation_p = citation.text_frame.paragraphs[0]
            citation_p.text = "Source: " + "; ".join(item.get("label", "") for item in citations)
            citation_p.font.size = PptPt(9); citation_p.font.italic = True; citation_p.font.color.rgb = PPT_SLATE
        safe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(9.6), PptInches(1.8), PptInches(2.75), PptInches(4.55))
        safe.fill.solid(); safe.fill.fore_color.rgb = PptRGBColor(240, 244, 247)
        safe.line.color.rgb = PptRGBColor(200, 210, 218)
        safe.text_frame.paragraphs[0].text = "Avatar-safe area"
        safe.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        safe.text_frame.paragraphs[0].font.size = PptPt(11)
        safe.text_frame.paragraphs[0].font.color.rgb = PPT_SLATE
        _add_footer(slide, project, index + 1)
    if images:
        slide = deck.slides.add_slide(blank)
        _add_slide_title(slide, "Supporting visuals", "Notebook and author assets")
        item = images[0]
        slide.shapes.add_picture(str(item["path"]), PptInches(0.9), PptInches(1.75), width=PptInches(8.8), height=PptInches(4.95))
        caption = slide.shapes.add_textbox(PptInches(10), PptInches(2.0), PptInches(2.4), PptInches(2.0))
        caption_p = caption.text_frame.paragraphs[0]
        caption_p.text = item.get("caption") or item.get("alt_text") or ""
        caption_p.font.size = PptPt(14); caption_p.font.color.rgb = PPT_SLATE
        _add_footer(slide, project, len(deck.slides))
    path = output_dir / f"{_safe_slug(project['title'])}.pptx"
    deck.save(path)
    return path


def _render_quiz(project: dict, output_dir: pathlib.Path) -> list[pathlib.Path]:
    items = project.get("quiz_items") or []
    if not items:
        raise ValueError("The approved project does not contain generated quiz items.")
    data_path = output_dir / f"{_safe_slug(project['title'])}-quiz.json"
    data_path.write_text(json.dumps({"title": project["title"], "items": items}, indent=2), encoding="utf-8")
    sections = []
    for index, item in enumerate(items, start=1):
        options = "".join(f"<li>{html.escape(str(option))}</li>" for option in item.get("options", []))
        sections.append(f"<section><h2>{index}. {html.escape(item.get('question', ''))}</h2><ol>{options}</ol></section>")
    page = f"<!doctype html><html><head><meta charset='utf-8'><title>{html.escape(project['title'])}</title></head><body><main><h1>{html.escape(project['title'])}</h1>{''.join(sections)}</main></body></html>"
    html_path = output_dir / f"{_safe_slug(project['title'])}-quiz.html"
    html_path.write_text(page, encoding="utf-8")
    return [html_path, data_path]


def render_project(project: dict, brief: dict, template: dict, artifact_root: pathlib.Path) -> dict:
    render_id = str(uuid.uuid4())
    output_dir = artifact_root / render_id
    output_dir.mkdir(parents=True, exist_ok=False)
    images = _decode_images(project, output_dir)
    files: list[tuple[pathlib.Path, str]] = []
    fmt = project.get("format")
    if fmt == "job_aid":
        docx_path = _render_docx(project, brief, output_dir, images)
        files.append((docx_path, "DOCX"))
        pdf_path = _convert_pdf(docx_path, output_dir)
        if pdf_path:
            files.append((pdf_path, "PDF"))
    elif fmt == "presentation":
        files.append((_render_pptx(project, brief, output_dir, images), "PPTX"))
    elif fmt == "quiz":
        for path in _render_quiz(project, output_dir):
            files.append((path, path.suffix.lstrip(".").upper()))
    elif fmt == "video":
        raise ValueError("Video publishing requires the OpenAI narration and HeyGen scene worker; no local fallback is produced.")
    else:
        raise ValueError("Unsupported artifact format.")
    job = {
        "render_id": render_id,
        "project_id": project.get("project_id"),
        "status": "completed",
        "files": [_file_record(path, render_id, file_format) for path, file_format in files],
    }
    (output_dir / "render-job.json").write_text(json.dumps(job, indent=2), encoding="utf-8")
    return job
